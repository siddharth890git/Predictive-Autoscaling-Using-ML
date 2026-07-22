from pptx import Presentation

prs = Presentation('presentation/Predictive_Autoscaling_IBM_Presentation (1).pptx')
with open('ppt_text.txt', 'w', encoding='utf-8') as f:
    for i, slide in enumerate(prs.slides):
        f.write(f"\n--- Slide {i+1} ---\n")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                f.write(shape.text + "\n")
