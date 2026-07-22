import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# --- CONFIGURATION & COLORS ---
IBM_BLUE = RGBColor(15, 98, 254)
DARK_NAVY = RGBColor(5, 20, 56)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(244, 244, 244)
TEXT_DARK = RGBColor(22, 22, 22)
TEXT_MUTED = RGBColor(82, 82, 82)
CYAN = RGBColor(0, 114, 195)
GREEN = RGBColor(36, 161, 72)
ORANGE = RGBColor(241, 194, 27)

def set_slide_background(slide, color=DARK_NAVY):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_header_footer(slide, title_text, dark_mode=True):
    header_color = IBM_BLUE if dark_mode else LIGHT_GRAY
    text_color = WHITE if dark_mode else TEXT_DARK
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = title_text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.font.name = 'Segoe UI Semibold'
    
    line = slide.shapes.add_shape(9, Inches(0.5), Inches(6.8), Inches(12.33), Inches(0))
    line.line.color.rgb = IBM_BLUE if dark_mode else CYAN
    line.line.width = Pt(2)
    
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(5), Inches(0.4))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Predictive Autoscaling Using ML on Cloud | IBM Internship Project"
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE if dark_mode else TEXT_MUTED

def add_bullet_points(slide, content, left=0.5, top=1.5, width=12.33, height=5, dark_mode=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    text_color = WHITE if dark_mode else TEXT_DARK
    
    for i, point in enumerate(content):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = point
        p.font.size = Pt(20)
        p.font.name = 'Segoe UI'
        p.font.color.rgb = text_color
        p.space_after = Pt(14)
        if point.startswith("  -") or point.startswith("    -"):
            p.level = 1 if point.startswith("  -") else 2
            p.font.size = Pt(18)
            p.text = point.replace("- ", "").strip()

def create_slide_1_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_NAVY)
    
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Predictive Autoscaling Using\nMachine Learning on Cloud"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = 'Segoe UI'
    
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11), Inches(1))
    tf2 = txBox2.text_frame
    p2 = tf2.add_paragraph()
    p2.text = "IBM SkillsBuild Internship Final Review"
    p2.font.size = Pt(24)
    p2.font.color.rgb = CYAN
    
    txBox3 = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(5), Inches(1))
    tf3 = txBox3.text_frame
    p3 = tf3.add_paragraph()
    p3.text = "[Your Name]\n[Your University Name]\n[Date]"
    p3.font.size = Pt(16)
    p3.font.color.rgb = LIGHT_GRAY

def create_slide_2_agenda(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_NAVY)
    add_header_footer(slide, "Agenda")
    
    agenda_items = [
        "1. Introduction & Problem Statement",
        "2. Dataset & Feature Engineering",
        "3. Machine Learning Pipeline",
        "4. XGBoost Model Details",
        "5. Autoscaling Decision Engine",
        "6. System & Backend Architecture",
        "7. Dashboard Overview",
        "8. Results & Model Performance",
        "9. Current Implementation & Future Deployment",
        "10. Business Impact & Conclusion"
    ]
    add_bullet_points(slide, agenda_items, left=1, top=1.8, width=11, height=5)

def create_slide_3_introduction(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_NAVY)
    add_header_footer(slide, "Introduction")
    content = [
        "What is Autoscaling?",
        "  - The process of dynamically adjusting computational resources based on load.",
        "Reactive vs. Predictive Autoscaling",
        "  - Reactive: Scales only after traffic spikes happen (leads to latency & downtime).",
        "  - Predictive: Anticipates traffic spikes using ML and scales proactively.",
        "Project Goal",
        "  - Build a machine learning-driven autoscaler that predicts traffic 15 minutes ahead.",
        "  - Provide an intelligent decision engine to minimize SLA breaches and optimize costs."
    ]
    add_bullet_points(slide, content)

def create_slide_4_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_NAVY)
    add_header_footer(slide, "Problem Statement")
    content = [
        "Traditional Cloud Autoscaling is Reactive",
        "  - Rule-based triggers (e.g., CPU > 80%) have a boot-up delay.",
        "  - By the time instances are ready, user experience has already degraded.",
        "Resource Waste During Low Traffic",
        "  - Organizations over-provision servers to avoid downtime.",
        "  - Results in high infrastructure costs and inefficient resource utilization.",
        "The Need for Intelligence",
        "  - Cloud traffic exhibits seasonal and predictable patterns.",
        "  - We need a system that learns from historical data to scale preemptively."
    ]
    add_bullet_points(slide, content)

def create_slide_5_dataset(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_NAVY)
    add_header_footer(slide, "Dataset & Preprocessing")
    content = [
        "Base Dataset",
        "  - Comprehensive cloud traffic data with 100,000+ historical records.",
        "  - Recorded at 5-minute intervals covering CPU, Memory, IO, and Request Counts.",
        "Data Preprocessing",
        "  - Imputation of missing values and clamping of anomalies.",
        "  - Shifted the target variable by 3 steps to create a 15-minute forecasting horizon.",
        "Simulated Geographic Traffic Distribution",
        "  - Augmented base metrics with simulated demographic features.",
        "  - Includes Country, Device, Platform, and Subscription tiers to enable global analytics."
    ]
    add_bullet_points(slide, content)

def create_slide_6_feature_engineering(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_NAVY)
    add_header_footer(slide, "Feature Engineering")
    content = [
        "Extracted 38 distinct features to capture complex traffic patterns:",
        "Temporal Features",
        "  - Hour, Day, Month, Day of Week, Week, Quarter, Is_Weekend.",
        "Lag Features",
        "  - Request lags at 1, 3, 6, 12, and 24 steps to capture immediate trends.",
        "Rolling Window Statistics",
        "  - Rolling Mean and Standard Deviation over 6 and 12 periods.",
        "Derived Metrics",
        "  - Traffic Growth Rate, Request Velocity, Traffic Momentum.",
        "  - Resource Pressure and CPU/Memory Ratios."
    ]
    add_bullet_points(slide, content)

def create_slide_7_ml_pipeline(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_NAVY)
    add_header_footer(slide, "Machine Learning Pipeline")
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(4))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Historical Data → Preprocessing → Feature Engineering → XGBoost Training\n\n↓\n\nLive Traffic Simulator → Prediction Engine (+15m) → Autoscaling Rules\n\n↓\n\nReal-time Dashboard"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = CYAN
    p.font.name = 'Segoe UI'
    p.alignment = PP_ALIGN.CENTER

def create_slide_8_xgboost(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_NAVY)
    add_header_footer(slide, "XGBoost Model Implementation")
    content = [
        "Why XGBoost?",
        "  - Highly efficient gradient boosting framework perfect for structured tabular data.",
        "  - Selected over Prophet for the production API due to superior multivariate forecasting.",
        "Training Methodology",
        "  - 80/20 chronological split to preserve time-series integrity.",
        "  - Target variable shifted by 3 steps (predicting 15 minutes ahead).",
        "Model Hyperparameters",
        "  - n_estimators: 500",
        "  - learning_rate: 0.05",
        "  - max_depth: 8",
        "  - objective: reg:squarederror"
    ]
    add_bullet_points(slide, content)

def create_slide_9_autoscaler(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_NAVY)
    add_header_footer(slide, "Autoscaling Decision Engine")
    content = [
        "Rule-based deterministic engine evaluated against XGBoost predictions.",
        "Scale Up Triggers (Requires ANY condition):",
        "  - Capacity Utilization ≥ 85%",
        "  - CPU Usage ≥ 75%",
        "  - Memory Usage ≥ 75%",
        "Scale Down Triggers (Requires ALL conditions):",
        "  - Capacity Utilization ≤ 45%",
        "  - CPU Usage ≤ 30%",
        "  - Memory Usage ≤ 35%",
        "Maintain",
        "  - Triggers if neither Scale Up nor Scale Down thresholds are met.",
        "  - Instances are strictly clamped between a configured minimum and maximum."
    ]
    add_bullet_points(slide, content)

def create_slide_10_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_NAVY)
    add_header_footer(slide, "System Architecture")
    content = [
        "1. Traffic Simulator",
        "  - Iterates through historical test data to simulate real-world events.",
        "2. Prediction Layer",
        "  - XGBoost model ingests 38 features and outputs request count forecast.",
        "3. Decision Engine",
        "  - Autoscaler evaluates forecast against current instances and utilization.",
        "4. Backend API",
        "  - FastAPI serves the state, metrics, and history to the frontend.",
        "5. Frontend UI",
        "  - HTML/JS dashboard consumes REST APIs for visualization."
    ]
    add_bullet_points(slide, content)

def create_slide_11_backend(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_NAVY)
    add_header_footer(slide, "Backend API Architecture")
    content = [
        "High-performance FastAPI Server (`main.py`)",
        "Implemented REST Endpoints:",
        "  - GET / : Root status and version verification.",
        "  - GET /api/health : Infrastructure and sub-system health checks.",
        "  - GET /api/dashboard : Core unified payload delivering metrics, predictions,",
        "    autoscaler decisions, geo data, and alerts.",
        "  - GET /api/history : Full 50-step sliding window history for timeline charts.",
        "  - GET /api/analytics : Global aggregates and event counters.",
        "Live Event Loop",
        "  - Every API request advances the simulation, generating real-time predictions."
    ]
    add_bullet_points(slide, content)

def create_slide_12_dashboard_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_NAVY)
    add_header_footer(slide, "Dashboard Overview")
    content = [
        "Modern Responsive Frontend",
        "  - Built with Vanilla JS, HTML5, Bootstrap 5, and ApexCharts.",
        "  - Operates independently, decoupling UI from the Python backend.",
        "Key Panels",
        "  - KPI Grid: 8 live metrics including Traffic, Server Count, and CPU.",
        "  - Live Traffic: Visual comparison of Actual vs. Predicted requests.",
        "  - ML Prediction: Real-time error bars and model status.",
        "  - Autoscaling: Gantt-style timeline tracking Scale Up/Down decisions.",
        "  - Infrastructure: Overall system health score and alert logs."
    ]
    add_bullet_points(slide, content)

def create_slide_13_global_analytics(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_NAVY)
    add_header_footer(slide, "Global Analytics Dashboard")
    content = [
        "World Traffic Map",
        "  - Plotly-powered interactive choropleth visualizing request origin.",
        "Demographic Insights",
        "  - Live breakdown of Top Cloud Regions (e.g., us-east-1, ap-south-1).",
        "  - Distribution charts for Device Types and Operating Systems.",
        "  - Subscription and Content Type analytics.",
        "Interactive Features",
        "  - Country click-handlers displaying regional CPU and Memory usage.",
        "  - Data is heavily augmented by the backend's geo-simulation engine."
    ]
    add_bullet_points(slide, content)

def create_image_slide(prs, title, image_file, caption):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_NAVY)
    add_header_footer(slide, title)
    
    if os.path.exists(image_file) and os.path.getsize(image_file) > 0:
        try:
            slide.shapes.add_picture(image_file, Inches(2), Inches(1.5), Inches(9.33))
        except Exception:
            slide.shapes.add_shape(1, Inches(2), Inches(1.5), Inches(9.33), Inches(5))
    else:
        slide.shapes.add_shape(1, Inches(2), Inches(1.5), Inches(9.33), Inches(5))
        
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(12.33), Inches(0.5))
    txBox.text_frame.text = caption
    txBox.text_frame.paragraphs[0].font.color.rgb = CYAN
    txBox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    txBox.text_frame.paragraphs[0].font.size = Pt(16)

def create_slide_14_results_images(prs):
    create_image_slide(prs, "Traffic Trend Analysis", "img/cloud traffic trend over time.png", "The XGBoost model successfully captures traffic spikes for preemptive scaling.")
    create_image_slide(prs, "Hourly Traffic Distribution", "img/Average traffic by hour.png", "Analyzing request counts across hours to establish baseline utilization.")
    create_image_slide(prs, "Infrastructure Utilization", "img/CPU and memory utilization overtime.png", "Continuous monitoring of CPU and Memory metrics.")
    create_image_slide(prs, "Exploratory Data Analysis", "img/Correlation heat map.png", "Correlation Heatmap identifying linear and non-linear relationships.")

def create_slide_15_current_impl(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_NAVY)
    add_header_footer(slide, "Current Implementation")
    content = [
        "Containerization",
        "  - The project is fully containerized using Docker and Docker Compose.",
        "  - Ensures consistent environments across development and testing.",
        "Dashboard Implementations",
        "  - High-performance FastAPI backend paired with the HTML/JS frontend.",
        "  - Docker configuration currently targets the Streamlit application.",
        "  - FastAPI deployment configuration can be extended in future versions.",
        "Simulation Accuracy",
        "  - The system proves the viability of ML-driven scaling using strict logic",
        "    and historical datasets."
    ]
    add_bullet_points(slide, content)

def create_slide_16_future_deployment(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_NAVY)
    add_header_footer(slide, "Future Cloud Deployment")
    content = [
        "Transitioning from Simulation to Production:",
        "Amazon Web Services (AWS)",
        "  - Deploying the FastAPI backend to EC2 or Elastic Beanstalk.",
        "  - Replacing the Python autoscaler logic with AWS Auto Scaling Groups (ASG).",
        "  - Using AWS CloudWatch APIs to ingest real-time live traffic metrics.",
        "IBM Cloud Integration",
        "  - Deploying containers to IBM Cloud Kubernetes Service (IKS).",
        "  - Utilizing Watson Machine Learning for deploying the XGBoost model.",
        "  - Continuous monitoring via IBM Instana."
    ]
    add_bullet_points(slide, content)

def create_slide_17_business_impact(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_NAVY)
    add_header_footer(slide, "Business Impact")
    content = [
        "1. Elimination of SLA Breaches",
        "  - Pre-warming instances 15 minutes before traffic spikes guarantees zero",
        "    latency degradation for end users.",
        "2. Significant Cost Optimization",
        "  - Preemptively scaling down during quiet periods eliminates \"cloud waste\".",
        "  - Organizations pay only for the exact compute they need.",
        "3. Operational Efficiency",
        "  - Reduces the need for manual capacity planning and intervention.",
        "  - System automatically adapts to shifting global traffic patterns."
    ]
    add_bullet_points(slide, content)

def create_slide_18_limitations(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_NAVY)
    add_header_footer(slide, "Current Implementation Limitations")
    content = [
        "While the core architecture is robust, the current version has known limits:",
        "Demonstration Metrics",
        "  - Current dashboard metrics are demonstration values and should be replaced",
        "    with dynamically computed evaluation metrics in a production deployment.",
        "Simulated Geographic Traffic Distribution",
        "  - The geographic traffic distribution is simulated to prove the analytics UI.",
        "  - Future iterations must integrate with real CDN edge logs.",
        "Forecast Horizon",
        "  - The +15 minute forecast requires the boot time of the target cloud VMs to",
        "    be strictly under 15 minutes."
    ]
    add_bullet_points(slide, content)

def create_slide_19_conclusion(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_NAVY)
    add_header_footer(slide, "Conclusion")
    content = [
        "Predictive autoscaling is a critical requirement for modern cloud architectures.",
        "By leveraging XGBoost, we successfully shifted from a reactive penalty model",
        "to a proactive capacity model.",
        "",
        "The project demonstrates a complete end-to-end pipeline:",
        "  - Ingesting telemetry data.",
        "  - Generating accurate time-series forecasts.",
        "  - Applying deterministic rules to execute scaling actions.",
        "  - Visualizing the entire workflow on a modern, decoupled dashboard."
    ]
    add_bullet_points(slide, content)

def create_slide_20_thank_you(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_NAVY)
    txBox = slide.shapes.add_textbox(Inches(3), Inches(3), Inches(7.33), Inches(2))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Thank You\nAny Questions?"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = 'Segoe UI'
    p.alignment = PP_ALIGN.CENTER

def generate_presentation():
    print("Generating Detailed IBM-Style PPTX...")
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    create_slide_1_title(prs)
    create_slide_2_agenda(prs)
    create_slide_3_introduction(prs)
    create_slide_4_problem(prs)
    create_slide_5_dataset(prs)
    create_slide_6_feature_engineering(prs)
    create_slide_7_ml_pipeline(prs)
    create_slide_8_xgboost(prs)
    create_slide_9_autoscaler(prs)
    create_slide_10_architecture(prs)
    create_slide_11_backend(prs)
    create_slide_12_dashboard_overview(prs)
    create_slide_13_global_analytics(prs)
    create_slide_14_results_images(prs) # Creates 4 image slides
    create_slide_15_current_impl(prs)
    create_slide_16_future_deployment(prs)
    create_slide_17_business_impact(prs)
    create_slide_18_limitations(prs)
    create_slide_19_conclusion(prs)
    create_slide_20_thank_you(prs)
    
    save_path = os.path.join("presentation", "Predictive_Autoscaling_IBM_Presentation_Final.pptx")
    prs.save(save_path)
    print(f"✅ Presentation saved successfully to {save_path}")

if __name__ == "__main__":
    generate_presentation()
