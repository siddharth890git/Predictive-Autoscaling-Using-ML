import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# --- TEMPLATE COLORS ---
BG_COLOR = RGBColor(244, 246, 249)      # Light grayish blue
NAVY = RGBColor(4, 13, 54)              # Dark Navy
BRIGHT_GREEN = RGBColor(70, 209, 31)    # Vibrant Green
TEAL = RGBColor(0, 114, 112)            # Deep Teal
ACCENT_BLUE = RGBColor(15, 98, 254)     # IBM Blue
WHITE = RGBColor(255, 255, 255)
DARK_TEXT = RGBColor(30, 30, 30)

def set_background(slide, color=BG_COLOR):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title(slide, text):
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.8))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.font.name = 'Segoe UI Semibold'

def create_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, NAVY)
    
    # Add a decorative chevron
    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(8.5), Inches(0), Inches(5), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = BRIGHT_GREEN
    shape.line.fill.background()
    
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(7), Inches(2))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Predictive Autoscaling"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = 'Segoe UI'
    
    p2 = tf.add_paragraph()
    p2.text = "Machine Learning on Cloud"
    p2.font.size = Pt(36)
    p2.font.color.rgb = BRIGHT_GREEN
    
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(7), Inches(1))
    tf2 = txBox2.text_frame
    p3 = tf2.add_paragraph()
    p3.text = "IBM SkillsBuild Internship Final Review\n[Your Name] | [Your University]"
    p3.font.size = Pt(18)
    p3.font.color.rgb = WHITE

def create_chevron_process_slide(prs, title, steps):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG_COLOR)
    add_title(slide, title)
    
    num_steps = len(steps)
    width = 11.5 / num_steps
    colors = [BRIGHT_GREEN, NAVY, TEAL, ACCENT_BLUE, NAVY, BRIGHT_GREEN, TEAL]
    
    for i, step in enumerate(steps):
        left = 0.5 + (i * width)
        # Chevron Shape
        shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(left), Inches(2), Inches(width - 0.1), Inches(1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors[i % len(colors)]
        shape.line.fill.background()
        
        # Chevron Text
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = step['title']
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE
        
        # Details Box below Chevron
        txBox = slide.shapes.add_textbox(Inches(left), Inches(3.2), Inches(width - 0.1), Inches(3))
        tf_desc = txBox.text_frame
        tf_desc.word_wrap = True
        p_desc = tf_desc.add_paragraph()
        p_desc.text = step['desc']
        p_desc.font.size = Pt(14)
        p_desc.font.color.rgb = DARK_TEXT

def create_card_grid_slide(prs, title, cards, columns=3):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG_COLOR)
    add_title(slide, title)
    
    colors = [BRIGHT_GREEN, TEAL, NAVY, ACCENT_BLUE, BRIGHT_GREEN, TEAL, NAVY, ACCENT_BLUE]
    
    margin_left = 0.5
    margin_top = 1.8
    card_width = 12.33 / columns - 0.2
    card_height = 2.2
    
    for i, card in enumerate(cards):
        col = i % columns
        row = i // columns
        left = margin_left + (col * (card_width + 0.2))
        top = margin_top + (row * (card_height + 0.3))
        
        # Card Background
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(card_width), Inches(card_height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = colors[i % len(colors)]
        shape.line.width = Pt(2)
        
        # Decorator Badge
        badge = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(0.4), Inches(card_height))
        badge.fill.solid()
        badge.fill.fore_color.rgb = colors[i % len(colors)]
        badge.line.fill.background()
        
        # Card Text
        txBox = slide.shapes.add_textbox(Inches(left + 0.5), Inches(top + 0.1), Inches(card_width - 0.6), Inches(card_height - 0.2))
        tf = txBox.text_frame
        tf.word_wrap = True
        
        p = tf.add_paragraph()
        p.text = f"0{i+1}" if i < 9 else str(i+1)
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = colors[i % len(colors)]
        
        p2 = tf.add_paragraph()
        p2.text = card['title']
        p2.font.size = Pt(18)
        p2.font.bold = True
        p2.font.color.rgb = NAVY
        
        p3 = tf.add_paragraph()
        p3.text = card['desc']
        p3.font.size = Pt(14)
        p3.font.color.rgb = DARK_TEXT

def create_image_placeholder_slide(prs, title, image_file, desc_title, desc_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, BG_COLOR)
    add_title(slide, title)
    
    # Content Card
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(4), Inches(5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()
    
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.add_paragraph()
    p.text = desc_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = BRIGHT_GREEN
    
    p2 = tf.add_paragraph()
    p2.text = "\n" + desc_text
    p2.font.size = Pt(16)
    p2.font.color.rgb = WHITE
    
    # Image or Placeholder
    if os.path.exists(image_file) and os.path.getsize(image_file) > 0:
        try:
            slide.shapes.add_picture(image_file, Inches(5), Inches(1.5), Inches(7.5))
        except Exception:
            ph = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(1.5), Inches(7.5), Inches(5))
            ph.fill.solid()
            ph.fill.fore_color.rgb = WHITE
            ph.text_frame.text = f"Paste {image_file} Here"
    else:
        ph = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(1.5), Inches(7.5), Inches(5))
        ph.fill.solid()
        ph.fill.fore_color.rgb = WHITE
        ph.text_frame.text = f"Paste your screenshot here\n({os.path.basename(image_file)})"
        ph.text_frame.paragraphs[0].font.color.rgb = TEAL

def generate_presentation():
    print("Generating Infographic Style PPTX (21 Slides)...")
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # 1. Title
    create_title_slide(prs)
    
    # 2. Agenda (Cards)
    create_card_grid_slide(prs, "Agenda", [
        {"title": "Introduction", "desc": "Problem Statement & Project Goals"},
        {"title": "Data Pipeline", "desc": "Dataset & Feature Engineering"},
        {"title": "Machine Learning", "desc": "XGBoost Model & Training"},
        {"title": "System Architecture", "desc": "Backend API & Autoscaling Engine"},
        {"title": "Visual Analytics", "desc": "Dashboard & Global Insights"},
        {"title": "Business Impact", "desc": "Results & Future Deployment Roadmap"}
    ], columns=3)
    
    # 3. Introduction (Cards)
    create_card_grid_slide(prs, "Introduction to Autoscaling", [
        {"title": "What is it?", "desc": "The process of dynamically adjusting computational resources based on load."},
        {"title": "Reactive Scaling", "desc": "Scales only after traffic spikes happen, leading to latency and downtime."},
        {"title": "Predictive Scaling", "desc": "Anticipates traffic using ML to scale proactively and prevent SLA breaches."},
        {"title": "Project Goal", "desc": "Build an ML-driven autoscaler that predicts traffic 15 minutes ahead."}
    ], columns=2)
    
    # 4. Problem Statement (Cards)
    create_card_grid_slide(prs, "Problem Statement", [
        {"title": "Reactive Autoscaling", "desc": "Rule-based triggers (CPU > 80%) have boot-up delays leading to bad UX."},
        {"title": "Resource Waste", "desc": "Over-provisioning servers during low traffic spikes costs."},
        {"title": "Missing Intelligence", "desc": "Current systems ignore historical seasonal traffic patterns."}
    ], columns=3)

    # 5. Dataset (Cards)
    create_card_grid_slide(prs, "Dataset & Preprocessing", [
        {"title": "Base Dataset", "desc": "100,000+ historical records at 5-minute intervals (CPU, Mem, Request Counts)."},
        {"title": "Data Preprocessing", "desc": "Imputed missing values, clamped anomalies, and shifted target by 3 steps (+15 min)."},
        {"title": "Simulated Geo Data", "desc": "Augmented with simulated demographic features (Country, Device) for analytics."}
    ], columns=3)
    
    # 6. Feature Engineering (Cards)
    create_card_grid_slide(prs, "Feature Engineering Strategies", [
        {"title": "Temporal Features", "desc": "Extracted Hour, Day, Month, Weekend metrics."},
        {"title": "Lag Features", "desc": "Request lags (1, 3, 6, 12, 24 steps) to capture immediate trends."},
        {"title": "Rolling Windows", "desc": "Rolling Mean & Std Dev over 6 and 12 periods."},
        {"title": "Derived Metrics", "desc": "Calculated Traffic Growth Rate & Velocity."},
        {"title": "Resource Ratios", "desc": "CPU/Memory Ratios to gauge instance pressure."},
        {"title": "Feature Count", "desc": "Total of 38 distinct features fed into the model."}
    ], columns=3)
    
    # 7. Machine Learning Pipeline (Chevron)
    create_chevron_process_slide(prs, "Machine Learning Pipeline", [
        {"title": "1. Dataset", "desc": "Historical records with CPU, Mem, IO"},
        {"title": "2. Feature Eng", "desc": "Created 38 features (Lags, Ratios)"},
        {"title": "3. XGBoost", "desc": "Forecasts 15-minutes (3 steps) ahead"},
        {"title": "4. Decision Engine", "desc": "Scale rules against forecast"},
        {"title": "5. Dashboard", "desc": "Live tracking via FastAPI"}
    ])

    # 8. XGBoost Model Details (Cards)
    create_card_grid_slide(prs, "XGBoost Model Implementation", [
        {"title": "Why XGBoost?", "desc": "Highly efficient gradient boosting framework perfect for structured tabular data."},
        {"title": "Training Setup", "desc": "80/20 chronological split to preserve time-series integrity."},
        {"title": "Hyperparameters", "desc": "n_estimators=500, learning_rate=0.05, max_depth=8."},
        {"title": "Objective", "desc": "reg:squarederror for regression forecasting."}
    ], columns=2)
    
    # 9. Autoscaling Rules (Cards)
    create_card_grid_slide(prs, "Autoscaling Decision Engine", [
        {"title": "Scale Up Logic", "desc": "ANY Condition Met:\n• Utilization ≥ 85%\n• CPU ≥ 75%\n• Memory ≥ 75%"},
        {"title": "Scale Down Logic", "desc": "ALL Conditions Met:\n• Utilization ≤ 45%\n• CPU ≤ 30%\n• Memory ≤ 35%"},
        {"title": "Maintain Logic", "desc": "If neither thresholds are met, the instance count remains steady and clamped to limits."}
    ], columns=3)
    
    # 10. System Architecture (Chevron)
    create_chevron_process_slide(prs, "System Architecture", [
        {"title": "1. Simulator", "desc": "Iterates historical data"},
        {"title": "2. Prediction Layer", "desc": "XGBoost ingest & forecast"},
        {"title": "3. Decision Engine", "desc": "Evaluates scaling rules"},
        {"title": "4. Backend API", "desc": "FastAPI serves state"},
        {"title": "5. Frontend UI", "desc": "HTML/JS visualizes data"}
    ])

    # 11. Backend API Architecture (Cards)
    create_card_grid_slide(prs, "Backend Architecture", [
        {"title": "FastAPI Server", "desc": "High performance REST API serving state and history."},
        {"title": "Endpoints", "desc": "Includes /api/dashboard, /api/history, and /api/analytics."},
        {"title": "Event Loop", "desc": "Every API request advances the simulation, generating real-time predictions."}
    ], columns=3)

    # 12. Dashboard Overview (Image)
    create_image_placeholder_slide(
        prs, "Dashboard Overview", "img/dashboard_screenshot.png",
        "Modern Analytics UI",
        "Built with HTML5, Bootstrap 5, and ApexCharts. \n\nOperates independently, decoupling UI from the Python backend with a live polling architecture."
    )

    # 13. Global Analytics Dashboard (Cards)
    create_card_grid_slide(prs, "Global Analytics Features", [
        {"title": "World Traffic Map", "desc": "Interactive choropleth visualizing request origin."},
        {"title": "Demographic Insights", "desc": "Breakdown of Top Cloud Regions (e.g., us-east-1)."},
        {"title": "Interactive Features", "desc": "Country click-handlers displaying regional CPU and Memory usage."}
    ], columns=3)

    # 14. Results - XGBoost (Image)
    create_image_placeholder_slide(
        prs, "Traffic Trend Analysis", "img/cloud traffic trend over time.png",
        "High Accuracy Forecasting",
        "The system successfully captures traffic spikes, allowing the autoscaler to pre-warm instances 15 minutes before actual load arrives."
    )
    
    # 15. Results - Traffic by Hour (Image)
    create_image_placeholder_slide(
        prs, "Hourly Traffic Distribution", "img/Average traffic by hour.png",
        "Seasonal Traffic Patterns",
        "Analyzing request counts across hours to establish baseline utilization and identify recurring peak periods."
    )

    # 16. Results - Utilization (Image)
    create_image_placeholder_slide(
        prs, "Infrastructure Utilization", "img/CPU and memory utilization overtime.png",
        "Resource Tracking",
        "Continuous monitoring of CPU and Memory metrics to ensure the Autoscaler triggers correctly before capacity limits are hit."
    )
    
    # 17. Results - Heatmap (Image)
    create_image_placeholder_slide(
        prs, "Exploratory Data Analysis", "img/Correlation heat map.png",
        "Correlation Heatmap",
        "Identifying linear and non-linear relationships between server metrics to select the most impactful features for the model."
    )
    
    # 18. Implementation & Future (Chevron)
    create_chevron_process_slide(prs, "Deployment Roadmap", [
        {"title": "1. Current State", "desc": "Dockerized Streamlit / FastAPI simulation engine."},
        {"title": "2. Transition", "desc": "Move FastAPI to AWS EC2 / Beanstalk."},
        {"title": "3. Cloud Integ", "desc": "Replace Python loop with AWS Auto Scaling Groups (ASG)."},
        {"title": "4. Production", "desc": "Continuous monitoring via CloudWatch and IBM Instana."}
    ])
    
    # 19. Limitations & Impact (Cards)
    create_card_grid_slide(prs, "Business Impact & Considerations", [
        {"title": "Eliminate SLA Breaches", "desc": "Pre-warming instances guarantees zero latency degradation for end users."},
        {"title": "Cost Optimization", "desc": "Scaling down preemptively eliminates costly cloud waste."},
        {"title": "Implementation Note", "desc": "Current dashboard metrics are demonstration values to be replaced in production."},
        {"title": "Geographic Simulation", "desc": "Geo-traffic is currently simulated; future versions will ingest real CDN logs."}
    ], columns=2)
    
    # 20. Conclusion (Cards)
    create_card_grid_slide(prs, "Conclusion", [
        {"title": "Proactive Capacity", "desc": "Shifted from a reactive penalty model to a proactive capacity model using XGBoost."},
        {"title": "End-to-End Flow", "desc": "Successfully ingested telemetry, generated forecasts, and applied scaling logic."},
        {"title": "Visual Delivery", "desc": "Visualized the entire workflow on a modern, decoupled real-time dashboard."}
    ], columns=3)
    
    # 21. Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide, NAVY)
    txBox = slide.shapes.add_textbox(Inches(3), Inches(3), Inches(7.33), Inches(2))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = "Thank You\nAny Questions?"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = 'Segoe UI'
    p.alignment = PP_ALIGN.CENTER
    
    # Save
    save_path = os.path.join("presentation", "Predictive_Autoscaling_IBM_Infographic_PPT.pptx")
    prs.save(save_path)
    print(f"✅ Presentation saved to {save_path}")

if __name__ == "__main__":
    generate_presentation()
