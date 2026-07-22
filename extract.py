from pptx import Presentation

prs = Presentation('presentation/Predictive_Autoscaling_IBM_Presentation (1).pptx')
for i, slide in enumerate(prs.slides):
    title = slide.shapes.title.text if slide.shapes.title else "No Title"
    print(f"Slide {i+1}: {title}")
